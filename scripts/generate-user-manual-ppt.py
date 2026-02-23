"""台股智慧選股系統 - 使用者操作手冊 PPT（圖示版）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Theme ---
BG = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x25, 0x25, 0x40)
GOLD = RGBColor(0xE5, 0xA9, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DIM = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x67, 0xC2, 0x3A)
BLUE = RGBColor(0x40, 0x9E, 0xFF)
RED = RGBColor(0xF5, 0x6C, 0x6C)
PURPLE = RGBColor(0xB3, 0x7F, 0xEB)
ORANGE = RGBColor(0xE6, 0xA2, 0x3C)
TEAL = RGBColor(0x13, 0xCE, 0xA0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BG


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Microsoft JhengHei"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def card(slide, l, t, w, h, fill=CARD, radius=0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def icon_card(slide, l, t, w, h, icon, title, desc, icon_color=GOLD):
    """Card with emoji icon, title, description."""
    card(slide, l, t, w, h)
    txt(slide, l + 0.15, t + 0.15, w - 0.3, 0.6, icon, sz=32, align=PP_ALIGN.CENTER)
    txt(slide, l + 0.15, t + 0.8, w - 0.3, 0.4, title, sz=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l + 0.15, t + 1.15, w - 0.3, 0.8, desc, sz=13, color=GRAY, align=PP_ALIGN.CENTER)


def bar(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()


def step_arrow(slide, l, t, w, h, label, color=GOLD):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft JhengHei"
    p.alignment = PP_ALIGN.CENTER


# ========== Slide 1: 封面 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 1, 1.8, 11, 1.2, "台股智慧選股系統", sz=52, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
bar(s, 4.5, 3.2, 4)
txt(s, 1, 3.5, 11, 0.8, "使用者操作手冊", sz=28, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 1, 4.8, 11, 0.5, "多因子篩選  ×  AI 智慧分析  ×  動能信號偵測", sz=18, color=DIM, align=PP_ALIGN.CENTER)
txt(s, 1, 6.2, 11, 0.4, "v2.2  |  2026 年 2 月", sz=14, color=DIM, align=PP_ALIGN.CENTER)

# ========== Slide 2: 系統三大核心 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "系統三大核心", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

icon_card(s, 1, 1.6, 3.4, 2.2, "📊", "多因子評分", "結合籌碼面、基本面、技術面\n三維度量化評分 0-100 分")
icon_card(s, 5, 1.6, 3.4, 2.2, "🤖", "AI 智慧分析", "Google Gemini AI 自動分析\n新聞摘要、投資建議、風險提示")
icon_card(s, 9, 1.6, 3.4, 2.2, "📈", "動能信號", "6 個右側進場信號偵測\n買賣點預測與動作建議")

# Bottom row
icon_card(s, 1, 4.2, 3.4, 2.2, "⏰", "每日自動更新", "每日 16:30 自動收集數據\n評分排名即時更新")
icon_card(s, 5, 4.2, 3.4, 2.2, "🔍", "智慧搜尋", "支援股票代碼或名稱搜尋\n非 Pipeline 股票自動補充資料")
icon_card(s, 9, 4.2, 3.4, 2.2, "👤", "會員系統", "Free / Premium 兩種等級\n差異化配額與功能")

# ========== Slide 3: 帳號註冊與登入 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "帳號註冊與登入", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

# Flow steps
step_arrow(s, 1, 1.6, 2.5, 0.8, "1. 填寫資料", BLUE)
step_arrow(s, 3.8, 1.6, 2.5, 0.8, "2. 選擇方案", GREEN)
step_arrow(s, 6.6, 1.6, 2.5, 0.8, "3. 完成註冊", GOLD)
step_arrow(s, 9.4, 1.6, 2.5, 0.8, "4. 登入使用", PURPLE)

card(s, 1, 2.8, 5.3, 3.5)
txt(s, 1.3, 2.9, 4.8, 0.5, "📝  註冊流程", sz=22, color=GOLD, bold=True)
txt(s, 1.3, 3.4, 4.8, 0.4, "• 輸入使用者名稱（唯一）", sz=16, color=GRAY)
txt(s, 1.3, 3.8, 4.8, 0.4, "• 輸入電子郵件（唯一）", sz=16, color=GRAY)
txt(s, 1.3, 4.2, 4.8, 0.4, "• 設定密碼（至少 8 字元）", sz=16, color=GRAY)
txt(s, 1.3, 4.6, 4.8, 0.4, "• 選擇 Free 或 Premium 方案", sz=16, color=GRAY)
txt(s, 1.3, 5.0, 4.8, 0.5, "• 註冊成功 → 自動跳轉登入", sz=16, color=GRAY)

card(s, 7, 2.8, 5.3, 3.5)
txt(s, 7.3, 2.9, 4.8, 0.5, "🔑  登入方式", sz=22, color=GOLD, bold=True)
txt(s, 7.3, 3.4, 4.8, 0.4, "• 輸入帳號 + 密碼", sz=16, color=GRAY)
txt(s, 7.3, 3.8, 4.8, 0.4, "• 系統自動驗證身份", sz=16, color=GRAY)
txt(s, 7.3, 4.2, 4.8, 0.4, "• 登入有效期 24 小時", sz=16, color=GRAY)
txt(s, 7.3, 4.6, 4.8, 0.4, "• 過期後自動導回登入頁面", sz=16, color=GRAY)
txt(s, 7.3, 5.0, 4.8, 0.5, "• 側邊欄顯示會員等級徽章", sz=16, color=GRAY)

# ========== Slide 4: 主儀表板 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "主儀表板 Dashboard", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

icon_card(s, 0.8, 1.5, 2.8, 2, "📋", "統計摘要", "分析股票數、今日最高分\n高籌碼分數、資料日期")
icon_card(s, 4, 1.5, 2.8, 2, "🏷️", "產業分類", "全部 + 各產業標籤\n快速切換篩選")
icon_card(s, 7.2, 1.5, 2.8, 2, "🏆", "Top 30 排名", "9 欄可排序表格\n每頁 10 筆分頁")
icon_card(s, 10.4, 1.5, 2.8, 2, "🔗", "一鍵深入", "點擊任一股票\n進入詳情頁面")

card(s, 0.8, 4, 11.8, 2.5)
txt(s, 1.1, 4.1, 5, 0.4, "排名表格欄位", sz=20, color=GOLD, bold=True)

# Table header simulation – match dashboard-view.vue columns
cols = [("#", 0.7), ("代號", 1.1), ("名稱", 1.5), ("總分", 1.0), ("收盤價", 1.3), ("漲跌", 1.1), ("籌碼", 1.1), ("基本面", 1.2), ("技術面", 1.2)]
widths = [w for _, w in cols]
for i, (col, w) in enumerate(cols):
    x = 1.1 + sum(widths[:i])
    card(s, x, 4.7, w - 0.1, 0.5, fill=RGBColor(0x35, 0x35, 0x55))
    txt(s, x, 4.7, w - 0.1, 0.5, col, sz=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

txt(s, 1.1, 5.5, 11, 0.5, "支援多欄位排序  |  每頁 10 筆分頁  |  Top 30 統計摘要卡片", sz=15, color=DIM, align=PP_ALIGN.CENTER)

# ========== Slide 5: 個股詳情頁 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "個股詳情頁", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

# 4 sections
icon_card(s, 0.8, 1.5, 2.8, 2.3, "📋", "評分卡片", "籌碼 / 基本面 / 技術面\n各項子指標明細分數")
icon_card(s, 4, 1.5, 2.8, 2.3, "🕯️", "K 線圖", "含 5 條均線指標\nMA5 / 10 / 20 / 60 / 120")
icon_card(s, 7.2, 1.5, 2.8, 2.3, "📈", "技術指標", "KD / MACD / RSI\n布林通道等圖表")
icon_card(s, 10.4, 1.5, 2.8, 2.3, "🤖", "AI 分析報告", "新聞摘要 + 投資建議\n情緒分析 + 風險提示")

card(s, 0.8, 4.3, 5.6, 2.4)
txt(s, 1.1, 4.4, 5, 0.5, "🎯  右側買法信號", sz=20, color=GOLD, bold=True)
txt(s, 1.1, 4.9, 5, 0.35, "• 6 個動能進場信號即時偵測", sz=15, color=GRAY)
txt(s, 1.1, 5.25, 5, 0.35, "• 買賣點預測：進場 / 停損 / 目標價", sz=15, color=GRAY)
txt(s, 1.1, 5.6, 5, 0.35, "• 動作建議：Buy / Hold / Avoid", sz=15, color=GRAY)

card(s, 6.8, 4.3, 5.6, 2.4)
txt(s, 7.1, 4.4, 5, 0.5, "🔍  搜尋導航", sz=20, color=GOLD, bold=True)
txt(s, 7.1, 4.9, 5, 0.35, "• 搜尋欄輸入代碼或名稱即可導航", sz=15, color=GRAY)
txt(s, 7.1, 5.25, 5, 0.35, "• 非熱門股自動補充 6 個月歷史資料", sz=15, color=GRAY)
txt(s, 7.1, 5.6, 5, 0.35, "• 鍵盤 ↑↓ Enter 快速操作", sz=15, color=GRAY)

# ========== Slide 6: AI 報告生成 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "AI 報告生成", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

# 3 states flow
card(s, 1, 1.5, 3.4, 2.5, fill=RGBColor(0x1E, 0x3A, 0x5F))
txt(s, 1.2, 1.6, 3, 0.5, "🆕  無報告", sz=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.2, 2.2, 3, 0.4, "按鈕顯示", sz=14, color=DIM, align=PP_ALIGN.CENTER)
txt(s, 1.2, 2.6, 3, 0.5, "「產生 AI 分析」", sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.2, 3.2, 3, 0.4, "✅ 可點擊", sz=14, color=GREEN, align=PP_ALIGN.CENTER)

card(s, 5, 1.5, 3.4, 2.5, fill=RGBColor(0x3A, 0x3A, 0x1E))
txt(s, 5.2, 1.6, 3, 0.5, "⏰  超過 24h", sz=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 5.2, 2.2, 3, 0.4, "按鈕顯示", sz=14, color=DIM, align=PP_ALIGN.CENTER)
txt(s, 5.2, 2.6, 3, 0.5, "「更新分析」", sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 5.2, 3.2, 3, 0.4, "✅ 可點擊", sz=14, color=GREEN, align=PP_ALIGN.CENTER)

card(s, 9, 1.5, 3.4, 2.5, fill=RGBColor(0x1E, 0x3A, 0x1E))
txt(s, 9.2, 1.6, 3, 0.5, "✅  24h 內", sz=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 9.2, 2.2, 3, 0.4, "按鈕顯示", sz=14, color=DIM, align=PP_ALIGN.CENTER)
txt(s, 9.2, 2.6, 3, 0.5, "「今日已分析」", sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 9.2, 3.2, 3, 0.4, "🔒 已禁用", sz=14, color=DIM, align=PP_ALIGN.CENTER)

card(s, 1, 4.5, 11.4, 2.2)
txt(s, 1.3, 4.6, 10.8, 0.5, "📋  AI 報告內容", sz=22, color=GOLD, bold=True)

icon_card(s, 1.3, 5.1, 2.4, 1.4, "📰", "新聞摘要", "近 14 天\n相關新聞整理")
icon_card(s, 4, 5.1, 2.4, 1.4, "💡", "投資建議", "AI 分析\n操作建議")
icon_card(s, 6.7, 5.1, 2.4, 1.4, "😊", "情緒分析", "市場情緒\n正面/負面判斷")
icon_card(s, 9.4, 5.1, 2.4, 1.4, "⚠️", "風險提示", "潛在風險\n注意事項")

# ========== Slide 7: 篩選功能 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "篩選功能", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

card(s, 0.8, 1.5, 5.8, 5)
txt(s, 1.1, 1.6, 5.2, 0.5, "🎯  自訂篩選", sz=24, color=BLUE, bold=True)
txt(s, 1.1, 2.2, 5.2, 0.4, "自由調整三因子權重比例", sz=16, color=GRAY)
# Weight bars
for i, (label, pct, color) in enumerate([("籌碼面", 40, GOLD), ("基本面", 35, GREEN), ("技術面", 25, BLUE)]):
    y = 2.8 + i * 0.7
    txt(s, 1.3, y, 1.5, 0.35, label, sz=15, color=WHITE)
    # Bar background
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(y + 0.05), Inches(3), Inches(0.25)) if False else None
    card(s, 2.8, y + 0.05, 3, 0.25, fill=RGBColor(0x35, 0x35, 0x55))
    card(s, 2.8, y + 0.05, 3 * pct / 100, 0.25, fill=color)
    txt(s, 2.8, y, 3, 0.35, f"{pct}%", sz=13, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, 1.1, 5.0, 5.2, 0.4, "結果支援排序、分頁，點擊進入個股詳情", sz=14, color=DIM)

card(s, 7, 1.5, 5.8, 5)
txt(s, 7.3, 1.6, 5.2, 0.5, "📈  右側買法篩選", sz=24, color=GREEN, bold=True)
txt(s, 7.3, 2.2, 5.2, 0.4, "偵測 6 個動能進場信號", sz=16, color=GRAY)

signals = [("量價齊揚", 25), ("突破20日高點", 20), ("MACD黃金交叉", 20), ("站回MA20", 15), ("KD低檔黃金交叉", 12), ("突破布林上軌", 8)]
for i, (name, weight) in enumerate(signals):
    y = 2.8 + i * 0.38
    txt(s, 7.5, y, 3, 0.35, name, sz=13, color=GRAY)
    txt(s, 10.5, y, 1, 0.35, f"{weight}分", sz=13, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)

txt(s, 7.3, 5.3, 5.2, 0.4, "可設定「最少信號數」+ 4 種進階篩選", sz=14, color=DIM)

# ========== Slide 7.5: 右側買法進階篩選 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "右側買法：進階篩選條件", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)
txt(s, 0.8, 1.2, 11, 0.35, "除了最少信號數，還有 4 個進階篩選條件可組合使用", sz=16, color=DIM)

icon_card(s, 0.8, 1.8, 2.8, 2.3, "🚀", "今日突破", "量價齊揚 + 突破20日高\n多方突破確認")
icon_card(s, 4, 1.8, 2.8, 2.3, "📈", "週趨勢向上", "MA5 > MA20 且持續上升\n中短期多頭排列")
icon_card(s, 7.2, 1.8, 2.8, 2.3, "⚠️", "風險等級", "低 / 中 / 高 三級\n依波動率+分數判定")
icon_card(s, 10.4, 1.8, 2.8, 2.3, "🌟", "強力推薦", "分數≥60、觸發≥3\n趨勢向上、非高風險")

card(s, 0.8, 4.5, 11.8, 2.2)
txt(s, 1.1, 4.6, 11.2, 0.5, "💡  條件可組合疊加，逐步縮小篩選範圍", sz=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.1, 5.2, 11.2, 0.4, "篩選結果表格新增「評級」「條件標籤」「訊號明細」三個欄位", sz=16, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, 1.1, 5.7, 11.2, 0.4, "候選池：最新交易日成交量前 100 名（含連假資料不足自動回退）", sz=15, color=DIM, align=PP_ALIGN.CENTER)

# ========== Slide 8: AI 聊天助手 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "AI 聊天助手", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

# Chat simulation
card(s, 3, 1.5, 7, 4.5, fill=RGBColor(0x20, 0x20, 0x38))
txt(s, 3.3, 1.6, 6.4, 0.5, "🤖  AI 台股投資助手", sz=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# User message
card(s, 6, 2.3, 3.5, 0.6, fill=RGBColor(0x2C, 0x5F, 0x2C))
txt(s, 6.1, 2.3, 3.3, 0.6, "台積電最近表現如何？", sz=14, color=WHITE, align=PP_ALIGN.RIGHT)

# AI response
card(s, 3.5, 3.1, 5, 1.2, fill=RGBColor(0x30, 0x30, 0x50))
txt(s, 3.6, 3.1, 4.8, 1.2, "台積電 (2330) 近期籌碼面表現強勢，\n外資連續 5 日買超。技術面 MACD 呈現\n黃金交叉，建議關注 ...", sz=13, color=GRAY)

# User message 2
card(s, 6.5, 4.5, 3, 0.6, fill=RGBColor(0x2C, 0x5F, 0x2C))
txt(s, 6.6, 4.5, 2.8, 0.6, "推薦哪些標的？", sz=14, color=WHITE, align=PP_ALIGN.RIGHT)

txt(s, 3.3, 5.3, 6.4, 0.4, "右下角浮動氣泡  ➜  點擊展開聊天面板", sz=14, color=DIM, align=PP_ALIGN.CENTER)

# Features on right
card(s, 0.8, 1.5, 1.8, 4.5)
txt(s, 0.9, 1.6, 1.6, 0.4, "功能特色", sz=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
for i, feat in enumerate(["自由對話", "股票分析", "市場趨勢", "投資建議", "歷史記錄"]):
    txt(s, 0.9, 2.2 + i * 0.55, 1.6, 0.4, f"• {feat}", sz=13, color=GRAY)

# ========== Slide 9: 其他功能 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "更多功能", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

icon_card(s, 0.8, 1.5, 3.6, 2.2, "📊", "籌碼統計", "法人買賣超趨勢圖\n融資融券變化追蹤")
icon_card(s, 4.8, 1.5, 3.6, 2.2, "⏪", "歷史回測", "驗證篩選策略績效\n選擇歷史日期回溯評分")
icon_card(s, 8.8, 1.5, 3.6, 2.2, "⚙️", "系統設定", "調整三因子權重比例\n設定每日自動執行時間")

icon_card(s, 0.8, 4.2, 3.6, 2.2, "📄", "報告清單", "查看所有 AI 分析報告\n搜尋與瀏覽歷史報告")
icon_card(s, 4.8, 4.2, 3.6, 2.2, "👤", "會員資料", "查看會員等級與配額\n修改個人資料")
icon_card(s, 8.8, 4.2, 3.6, 2.2, "🛡️", "管理員", "使用者管理\n會員等級調整")

# ========== Slide 10: 會員方案比較 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "會員方案比較", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

# Free plan
card(s, 1, 1.5, 5.2, 5, fill=RGBColor(0x22, 0x22, 0x3D))
txt(s, 1.3, 1.7, 4.6, 0.6, "Free 免費方案", sz=28, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
bar(s, 2, 2.4, 3.2)
features_free = [
    ("✅", "基本篩選與排名功能"),
    ("✅", "個股詳情與圖表"),
    ("✅", "右側買法信號"),
    ("5️⃣", "AI 報告：每日 5 份"),
    ("💬", "AI 聊天：每分鐘 3 則 / 每日 10 則"),
    ("🔒", "僅能查看自己生成的報告"),
]
for i, (icon, text) in enumerate(features_free):
    txt(s, 1.5, 2.7 + i * 0.5, 0.4, 0.4, icon, sz=16, align=PP_ALIGN.CENTER)
    txt(s, 2.1, 2.7 + i * 0.5, 3.8, 0.4, text, sz=15, color=GRAY)

# Premium plan
card(s, 7, 1.5, 5.2, 5, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 7.3, 1.7, 4.6, 0.6, "Premium 進階方案", sz=28, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
bar(s, 8, 2.4, 3.2)
features_premium = [
    ("✅", "所有 Free 功能"),
    ("✅", "個股詳情與圖表"),
    ("✅", "右側買法信號"),
    ("♾️", "AI 報告：無限制"),
    ("💬", "AI 聊天：每分鐘 5 則 / 每日 100 則"),
    ("🌟", "可查看所有使用者的報告"),
]
for i, (icon, text) in enumerate(features_premium):
    txt(s, 7.5, 2.7 + i * 0.5, 0.4, 0.4, icon, sz=16, align=PP_ALIGN.CENTER)
    txt(s, 8.1, 2.7 + i * 0.5, 3.8, 0.4, text, sz=15, color=WHITE)

# ========== Slide 11: 日常使用流程 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 0.8, 0.4, 11, 0.7, "日常使用流程", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 2.5)

steps = [
    ("1", "登入", "開啟系統\n輸入帳密登入", BLUE),
    ("2", "看排名", "主儀表板\n查看最新排名", GREEN),
    ("3", "選個股", "點擊感興趣\n的股票深入", GOLD),
    ("4", "看分析", "K線+評分\n+AI 報告", PURPLE),
    ("5", "找機會", "右側買法\n篩選進場點", ORANGE),
    ("6", "問 AI", "聊天助手\n即時諮詢", TEAL),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.8 + i * 2.05
    card(s, x, 1.5, 1.85, 3.5, fill=RGBColor(0x25, 0x25, 0x40))
    # Number circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.55), Inches(1.7), Inches(0.7), Inches(0.7))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Microsoft JhengHei"

    txt(s, x + 0.1, 2.6, 1.65, 0.5, title, sz=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, 3.1, 1.65, 1, desc, sz=14, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, 0.8, 5.5, 11.4, 0.5, "💡 系統每日 16:30 自動更新數據，登入即可看到最新結果", sz=16, color=DIM, align=PP_ALIGN.CENTER)

# ========== Slide 12: 結尾 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, 1, 2.2, 11, 1, "開始使用", sz=48, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
bar(s, 4.5, 3.3, 4)
txt(s, 1, 3.6, 11, 0.8, "台股智慧選股系統 v2.2", sz=24, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 1, 4.8, 11, 0.5, "多因子篩選  ×  AI 智慧分析  ×  動能信號偵測", sz=18, color=DIM, align=PP_ALIGN.CENTER)

# Save
out = "/Users/weihung/Desktop/project/stock-system/docs/台股智慧選股系統-使用者操作手冊.pptx"
prs.save(out)
print(f"Generated: {out} ({len(prs.slides)} slides)")

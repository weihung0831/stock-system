"""台股智慧選股系統 - 指標公式與評分邏輯 PPT（使用者版，含完整公式）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Theme ---
BG = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x25, 0x25, 0x40)
GOLD = RGBColor(0xE5, 0xA9, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DIM = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x67, 0xC2, 0x3A)
BLUE = RGBColor(0x40, 0x9E, 0xFF)
RED = RGBColor(0xF5, 0x6C, 0x6C)
PURPLE = RGBColor(0xB3, 0x7F, 0xEB)
ORANGE = RGBColor(0xE6, 0xA2, 0x3C)
TEAL = RGBColor(0x13, 0xCE, 0xA0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg_fill(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BG


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Microsoft JhengHei"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def card(slide, l, t, w, h, fill=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def bar(slide, l, t, w, color=GOLD):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def pct_bar(slide, l, t, w_total, pct, label, color, value_text=""):
    card(slide, l, t, w_total, 0.35, fill=RGBColor(0x35, 0x35, 0x55))
    if pct > 0:
        card(slide, l, t, w_total * pct / 100, 0.35, fill=color)
    txt(slide, l - 1.5, t, 1.4, 0.35, label, sz=13, color=WHITE, align=PP_ALIGN.RIGHT)
    txt(slide, l + w_total + 0.1, t, 0.8, 0.35, value_text or f"{pct}%", sz=13, color=color, bold=True)


def circle_num(slide, x, y, num, color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = str(num)
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Microsoft JhengHei"


def score_table(slide, l, t, w, rows, title_color=GOLD):
    """Draw a scoring threshold table. rows = [(condition, score, color), ...]"""
    row_h = 0.32
    for i, (cond, score_val, clr) in enumerate(rows):
        y = t + i * row_h
        card(slide, l, y, w * 0.65, row_h - 0.02, fill=RGBColor(0x2A, 0x2A, 0x45))
        card(slide, l + w * 0.65 + 0.05, y, w * 0.30, row_h - 0.02, fill=clr)
        txt(slide, l + 0.08, y, w * 0.6, row_h, cond, sz=11, color=GRAY, font="Consolas")
        txt(slide, l + w * 0.65 + 0.08, y, w * 0.28, row_h, score_val, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ========== Slide 1: 封面 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 1, 1.8, 11, 1.2, "台股智慧選股系統", sz=52, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
bar(s, 4.5, 3.2, 4)
txt(s, 1, 3.5, 11, 0.8, "指標公式與評分邏輯", sz=28, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 1, 4.8, 11, 0.5, "完整公式對照手冊", sz=18, color=DIM, align=PP_ALIGN.CENTER)
txt(s, 1, 6.2, 11, 0.4, "v2.2  |  2026 年 2 月", sz=14, color=DIM, align=PP_ALIGN.CENTER)

# ========== Slide 2: 評分系統總覽 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "評分系統如何運作？", sz=36, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)

steps = [
    ("1", "篩選候選股", "從 ~1,300 檔中\n篩出約 500 檔", BLUE),
    ("2", "三因子評分", "籌碼+基本面+技術\n各打 0-100 分", GREEN),
    ("3", "加權合計", "依權重合計\n產出綜合分數\n寫入排名", GOLD),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = 1.5 + i * 3.8
    card(s, x, 1.4, 3.2, 2.8)
    circle_num(s, x + 1.25, 1.55, num, color)
    txt(s, x + 0.1, 2.3, 3.0, 0.5, title, sz=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, 2.8, 3.0, 0.8, desc, sz=14, color=GRAY, align=PP_ALIGN.CENTER)

for i in range(2):
    x = 4.8 + i * 3.8
    txt(s, x, 2.2, 0.5, 0.5, "→", sz=28, color=DIM, align=PP_ALIGN.CENTER)

card(s, 0.8, 4.7, 11.5, 2)
txt(s, 1.1, 4.8, 5, 0.5, "預設權重配置", sz=20, color=GOLD, bold=True)

pct_bar(s, 3.3, 5.4, 8, 40, "籌碼面", GOLD, "40%")
pct_bar(s, 3.3, 5.85, 8, 35, "基本面", GREEN, "35%")
pct_bar(s, 3.3, 6.3, 8, 25, "技術面", BLUE, "25%")

txt(s, 8, 4.8, 4, 0.5, "💡 可在設定頁面自訂權重", sz=14, color=DIM, align=PP_ALIGN.RIGHT)

# ========== Slide 3: 候選股篩選 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "Step 1  候選股篩選機制", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)

card(s, 0.8, 1.5, 5.5, 3.8, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 1.1, 1.6, 5, 0.5, "📊  路線 A：量能異常偵測", sz=22, color=BLUE, bold=True)
txt(s, 1.1, 2.2, 5, 0.35, "公式：", sz=14, color=DIM)
txt(s, 1.1, 2.5, 5, 0.4, "ratio = 本週成交量 ÷ 上週成交量", sz=16, color=WHITE, font="Consolas")
txt(s, 1.1, 2.9, 5, 0.4, "篩選條件：ratio > 2.5", sz=16, color=GREEN, bold=True, font="Consolas")
txt(s, 1.1, 3.5, 5, 0.8, "意義：成交量突然暴增的股票\n通常代表有重大事件或資金湧入", sz=13, color=GRAY)

card(s, 7, 1.5, 5.5, 3.8, fill=RGBColor(0x1E, 0x3A, 0x1E))
txt(s, 7.3, 1.6, 5, 0.5, "🏆  路線 B：Top 500 保底", sz=22, color=GREEN, bold=True)
txt(s, 7.3, 2.2, 5, 0.35, "篩選條件：", sz=14, color=DIM)
txt(s, 7.3, 2.5, 5, 0.4, "FALLBACK_TOP_N = 500", sz=16, color=WHITE, font="Consolas")
txt(s, 7.3, 2.9, 5, 0.4, "最新交易日成交量前 500 名", sz=16, color=GREEN, bold=True)
txt(s, 7.3, 3.5, 5, 0.8, "意義：確保市場上最活躍的\n股票都會被納入評分", sz=13, color=GRAY)

card(s, 3, 5.8, 7.3, 1.2, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 3.3, 5.85, 6.7, 0.5, "🎯  合併結果", sz=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
txt(s, 3.3, 6.35, 6.7, 0.5, "A 排前面 + B 填充補齊 → 去重 → 上限 500 檔候選股進入評分", sz=14, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, 3, 5.2, 3.5, 0.5, "↓", sz=28, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, 7, 5.2, 3.5, 0.5, "↓", sz=28, color=GREEN, align=PP_ALIGN.CENTER)

# ========== Slide 4: 籌碼面總覽 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "Step 2A  籌碼面評分 (chip_score)", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)
txt(s, 0.8, 1.15, 11, 0.4, "chip_score = A × 30% + B × 40% + C × 30%", sz=18, color=WHITE, font="Consolas")

# A: Consecutive buy
card(s, 0.8, 1.7, 3.8, 5.2, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 1.0, 1.8, 3.4, 0.5, "🏛️  A. 連續買超天數", sz=18, color=BLUE, bold=True)
txt(s, 1.0, 2.2, 3.4, 0.3, "佔比 30%", sz=13, color=DIM)

score_table(s, 1.0, 2.6, 3.4, [
    ("外資連買: 每天 +5 分", "最高 50", GREEN),
    ("投信連買: 每天 +3 分", "最高 30", BLUE),
    ("自營連買: 每天 +2 分", "最高 20", PURPLE),
])

txt(s, 1.0, 3.7, 3.4, 0.3, "total = 外資 + 投信 + 自營", sz=12, color=WHITE, font="Consolas")
txt(s, 1.0, 4.0, 3.4, 0.3, "score = min(total, 100)", sz=12, color=GOLD, font="Consolas")
txt(s, 1.0, 4.5, 3.4, 1.0, "💡 從最新交易日向前數\n只要出現賣超就中斷計數\n外資權重最高 (最高 50 分)", sz=11, color=DIM)

# B: Net buy ratio
card(s, 4.9, 1.7, 3.8, 5.2, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 5.1, 1.8, 3.4, 0.5, "📊  B. 法人買超佔比", sz=18, color=GREEN, bold=True)
txt(s, 5.1, 2.2, 3.4, 0.3, "佔比 40%（最重要）", sz=13, color=DIM)

txt(s, 5.1, 2.6, 3.4, 0.3, "公式：", sz=12, color=DIM)
txt(s, 5.1, 2.85, 3.4, 0.3, "ratio = 近5日法人淨買超", sz=11, color=WHITE, font="Consolas")
txt(s, 5.1, 3.1, 3.4, 0.3, "     ÷ 近5日成交量 × 100", sz=11, color=WHITE, font="Consolas")

score_table(s, 5.1, 3.5, 3.4, [
    ("ratio >= 5%", "100 分", GREEN),
    ("ratio = 0%", "50 分", GOLD),
    ("ratio <= -5%", "0 分", RED),
])

txt(s, 5.1, 4.6, 3.4, 0.3, "線性映射公式：", sz=12, color=DIM)
txt(s, 5.1, 4.9, 3.4, 0.3, "score = 50 + ratio × 5", sz=12, color=GOLD, bold=True, font="Consolas")
txt(s, 5.1, 5.3, 3.4, 0.5, "💡 夾在 0~100 之間\nmax(0, min(100, score))", sz=11, color=DIM)

# C: Margin
card(s, 9, 1.7, 3.8, 5.2, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 9.2, 1.8, 3.4, 0.5, "💰  C. 融資融券變化", sz=18, color=ORANGE, bold=True)
txt(s, 9.2, 2.2, 3.4, 0.3, "佔比 30%", sz=13, color=DIM)

txt(s, 9.2, 2.6, 3.4, 0.3, "5日變化率計算：", sz=12, color=DIM)
txt(s, 9.2, 2.85, 3.4, 0.3, "融資變化 = (最新-5日前)", sz=11, color=WHITE, font="Consolas")
txt(s, 9.2, 3.1, 3.4, 0.3, "        ÷ 5日前餘額 × 100", sz=11, color=WHITE, font="Consolas")

txt(s, 9.2, 3.5, 3.4, 0.3, "評分規則：基礎分 = 50", sz=12, color=DIM)

txt(s, 9.2, 3.85, 3.4, 0.25, "融資減少 → 加分", sz=12, color=GREEN)
txt(s, 9.2, 4.1, 3.4, 0.25, "  +min(|變化率|×2, 30)", sz=10, color=WHITE, font="Consolas")
txt(s, 9.2, 4.4, 3.4, 0.25, "融券增加 → 加分", sz=12, color=GREEN)
txt(s, 9.2, 4.65, 3.4, 0.25, "  +min(變化率×2, 20)", sz=10, color=WHITE, font="Consolas")

txt(s, 9.2, 5.0, 3.4, 0.3, "score = 50 + 融資 + 融券", sz=11, color=GOLD, bold=True, font="Consolas")
txt(s, 9.2, 5.35, 3.4, 0.8, "💡 融資減少=散戶退出（正面）\n融券增加=軋空潛力（正面）\n資料不足→中性 50 分", sz=11, color=DIM)

# ========== Slide 5: 基本面總覽 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "Step 2B  基本面評分 (fundamental_score)", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)

# Mode selection
txt(s, 0.8, 1.2, 11, 0.4, "三種評分模式（自動判斷資料可用性）", sz=18, color=WHITE)

# Mode 1: Full financial
card(s, 0.8, 1.8, 3.8, 5.2, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 1.0, 1.85, 3.4, 0.4, "📊  模式 A：有季財報", sz=17, color=BLUE, bold=True)
txt(s, 1.0, 2.2, 3.4, 0.3, "（7 指標加權）", sz=12, color=DIM)

indicators_a = [
    ("營收 YoY", "20%", GREEN),
    ("EPS 趨勢", "15%", BLUE),
    ("毛利率穩定度", "10%", TEAL),
    ("ROE", "15%", PURPLE),
    ("負債比", "15%", ORANGE),
    ("現金流", "15%", GREEN),
    ("本益比 PER", "10%", DIM),
]
for i, (name, w, clr) in enumerate(indicators_a):
    y = 2.6 + i * 0.38
    card(s, 1.0, y, 2.2, 0.32, fill=RGBColor(0x2A, 0x2A, 0x45))
    txt(s, 1.1, y, 2.0, 0.32, name, sz=11, color=GRAY)
    card(s, 3.25, y, 0.5, 0.32, fill=clr)
    txt(s, 3.25, y, 0.5, 0.32, w, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Mode 2: Revenue only
card(s, 4.9, 1.8, 3.8, 5.2, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 5.1, 1.85, 3.4, 0.4, "📈  模式 B：只有營收", sz=17, color=GREEN, bold=True)
txt(s, 5.1, 2.2, 3.4, 0.3, "（無季財報時自動切換）", sz=12, color=DIM)

indicators_b = [
    ("營收 YoY", "60%", GREEN),
    ("EPS 趨勢", "10%", BLUE),
    ("毛利率", "5%", TEAL),
    ("ROE", "5%", PURPLE),
    ("負債比", "5%", ORANGE),
    ("現金流", "5%", GREEN),
    ("本益比", "10%", DIM),
]
for i, (name, w, clr) in enumerate(indicators_b):
    y = 2.6 + i * 0.38
    card(s, 5.1, y, 2.2, 0.32, fill=RGBColor(0x2A, 0x2A, 0x45))
    txt(s, 5.2, y, 1.5, 0.32, name, sz=11, color=GRAY)
    card(s, 7.35, y, 0.55, 0.32, fill=clr)
    txt(s, 7.35, y, 0.55, 0.32, w, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, 5.1, 5.35, 3.4, 0.3, "💡 除營收外皆給中性 50 分", sz=11, color=DIM)

# Mode 3: Valuation
card(s, 9, 1.8, 3.8, 5.2, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 9.2, 1.85, 3.4, 0.4, "⚖️  模式 C：估值替代", sz=17, color=ORANGE, bold=True)
txt(s, 9.2, 2.2, 3.4, 0.3, "（完全無營收/財報時）", sz=12, color=DIM)

indicators_c = [
    ("PER 本益比", "30%", GOLD),
    ("PBR 股價淨值比", "30%", BLUE),
    ("殖利率", "40%", GREEN),
]
for i, (name, w, clr) in enumerate(indicators_c):
    y = 2.6 + i * 0.5
    card(s, 9.2, y, 2.4, 0.42, fill=RGBColor(0x2A, 0x2A, 0x45))
    txt(s, 9.3, y, 1.6, 0.42, name, sz=12, color=GRAY)
    card(s, 11.65, y, 0.55, 0.42, fill=clr)
    txt(s, 11.65, y, 0.55, 0.42, w, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, 9.2, 4.3, 3.4, 0.3, "💡 使用 TWSE 公開資料", sz=11, color=DIM)
txt(s, 9.2, 4.6, 3.4, 0.3, "total = PER×30%", sz=11, color=WHITE, font="Consolas")
txt(s, 9.2, 4.85, 3.4, 0.3, "     + PBR×30%", sz=11, color=WHITE, font="Consolas")
txt(s, 9.2, 5.1, 3.4, 0.3, "     + 殖利率×40%", sz=11, color=WHITE, font="Consolas")

# ========== Slide 6: 基本面公式細節 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "基本面：各指標評分門檻", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)

# Revenue YoY
card(s, 0.8, 1.3, 3.7, 3.0, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 1.0, 1.35, 3.3, 0.35, "📈 營收 YoY", sz=16, color=GREEN, bold=True)
score_table(s, 1.0, 1.75, 3.3, [
    ("YoY >= 20%", "100 分", GREEN),
    ("YoY >= 10%", "80 分", BLUE),
    ("YoY >= 5%", "60 分", TEAL),
    ("YoY >= 0%", "40 分", GOLD),
    ("YoY >= -5%", "20 分", ORANGE),
    ("YoY < -5%", "0 分", RED),
])

# EPS Trend
card(s, 4.8, 1.3, 3.7, 3.0, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 5.0, 1.35, 3.3, 0.35, "💹 EPS 趨勢", sz=16, color=BLUE, bold=True)
txt(s, 5.0, 1.75, 3.3, 0.3, "公式：", sz=12, color=DIM)
txt(s, 5.0, 2.0, 3.3, 0.3, "比較相鄰季度 EPS", sz=11, color=WHITE, font="Consolas")
txt(s, 5.0, 2.25, 3.3, 0.3, "ratio = 成長季數 ÷ 比較次數", sz=11, color=WHITE, font="Consolas")
txt(s, 5.0, 2.55, 3.3, 0.3, "score = ratio × 100", sz=12, color=GOLD, bold=True, font="Consolas")
txt(s, 5.0, 2.95, 3.3, 0.6, "例：4 季中 3 季成長\nscore = (3÷3) × 100 = 100\n不足 2 季 → 中性 50 分", sz=11, color=DIM)

# Gross Margin
card(s, 8.8, 1.3, 3.7, 3.0, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 9.0, 1.35, 3.3, 0.35, "📊 毛利率穩定度", sz=16, color=TEAL, bold=True)
txt(s, 9.0, 1.75, 3.3, 0.3, "兩個分數加權：", sz=12, color=DIM)
txt(s, 9.0, 2.05, 3.3, 0.3, "穩定分 = max(0, 100-標準差×10)", sz=10, color=WHITE, font="Consolas")
txt(s, 9.0, 2.3, 3.3, 0.3, "趨勢分 = 50 + (最新-最舊)×5", sz=10, color=WHITE, font="Consolas")
txt(s, 9.0, 2.65, 3.3, 0.3, "最終 = 穩定×60% + 趨勢×40%", sz=11, color=GOLD, bold=True, font="Consolas")
txt(s, 9.0, 3.05, 3.3, 0.5, "💡 波動越小越穩定\n趨勢上升則加分", sz=11, color=DIM)

# ROE
card(s, 0.8, 4.6, 3.7, 2.5, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 1.0, 4.65, 3.3, 0.35, "🏦 ROE 股東權益報酬率", sz=16, color=PURPLE, bold=True)
score_table(s, 1.0, 5.05, 3.3, [
    ("ROE >= 15%", "100 分", GREEN),
    ("ROE >= 12%", "80 分", BLUE),
    ("ROE >= 10%", "70 分", TEAL),
    ("ROE >= 8%", "60 分", GOLD),
    ("ROE >= 5%", "40 分", ORANGE),
    ("ROE >= 0%", "20 分", RED),
    ("ROE < 0%", "0 分", RED),
])

# Debt Ratio
card(s, 4.8, 4.6, 3.7, 2.5, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 5.0, 4.65, 3.3, 0.35, "⚖️ 負債比", sz=16, color=ORANGE, bold=True)
score_table(s, 5.0, 5.05, 3.3, [
    ("< 30%", "100 分", GREEN),
    ("< 50%", "80 分", BLUE),
    ("< 60%", "60 分", TEAL),
    ("< 70%", "40 分", GOLD),
    ("< 80%", "20 分", ORANGE),
    (">= 80%", "0 分", RED),
])

# Cash Flow + PE
card(s, 8.8, 4.6, 3.7, 2.5, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 9.0, 4.65, 3.3, 0.35, "💵 現金流 + 本益比", sz=16, color=GREEN, bold=True)

txt(s, 9.0, 5.05, 3.3, 0.25, "現金流：", sz=12, color=DIM)
score_table(s, 9.0, 5.3, 3.3, [
    ("營業現金流 > 0", "+50 分", GREEN),
    ("自由現金流 > 0", "+50 分", GREEN),
])

txt(s, 9.0, 5.98, 3.3, 0.25, "本益比 PER：", sz=12, color=DIM)
score_table(s, 9.0, 6.22, 3.3, [
    ("10~15 (理想)", "100 分", GREEN),
    ("<10 (偏低)", "80 分", BLUE),
    ("15~20 (偏高)", "70 分", GOLD),
    ("20~30 (昂貴)", "50 分", ORANGE),
    (">30 (過高)", "30 分", RED),
])

# ========== Slide 7: 估值替代評分 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "基本面：估值替代評分（模式 C 細節）", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)
txt(s, 0.8, 1.2, 11, 0.4, "當無營收也無財報時，使用 TWSE 公開的 PER/PBR/殖利率替代評分", sz=16, color=DIM)

# PER score
card(s, 0.8, 1.8, 3.8, 3.5, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 1.0, 1.85, 3.4, 0.4, "📈 PER 本益比 (30%)", sz=17, color=GOLD, bold=True)
score_table(s, 1.0, 2.35, 3.4, [
    ("PER <= 0 (虧損)", "50 分 (中性)", GOLD),
    ("PER < 10 (偏低)", "80 分", BLUE),
    ("10 <= PER < 15 (理想)", "100 分", GREEN),
    ("15 <= PER < 20", "70 分", TEAL),
    ("20 <= PER < 30", "50 分", ORANGE),
    ("PER >= 30", "30 分", RED),
])

# PBR score
card(s, 4.9, 1.8, 3.8, 3.5, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 5.1, 1.85, 3.4, 0.4, "📊 PBR 股價淨值比 (30%)", sz=17, color=BLUE, bold=True)
score_table(s, 5.1, 2.35, 3.4, [
    ("PBR <= 0", "50 分 (中性)", GOLD),
    ("PBR < 1.0", "90 分", GREEN),
    ("1.0 <= PBR < 1.5", "80 分", BLUE),
    ("1.5 <= PBR < 2.0", "60 分", TEAL),
    ("2.0 <= PBR < 3.0", "40 分", ORANGE),
    ("PBR >= 3.0", "20 分", RED),
])

# Dividend yield score
card(s, 9, 1.8, 3.8, 3.5, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 9.2, 1.85, 3.4, 0.4, "💰 殖利率 (40%)", sz=17, color=GREEN, bold=True)
score_table(s, 9.2, 2.35, 3.4, [
    ("殖利率 <= 0%", "30 分", RED),
    ("0% < 殖利率 < 2%", "30 分", RED),
    ("殖利率 >= 2%", "40 分", ORANGE),
    ("殖利率 >= 3%", "60 分", TEAL),
    ("殖利率 >= 4%", "80 分", BLUE),
    ("殖利率 >= 6%", "100 分", GREEN),
])

# Example
card(s, 0.8, 5.7, 11.7, 1.2, fill=RGBColor(0x20, 0x20, 0x38))
txt(s, 1.1, 5.75, 10, 0.4, "📝  範例：某股 PER=12, PBR=1.2, 殖利率=4.5%", sz=16, color=GOLD, bold=True)
txt(s, 1.1, 6.15, 10, 0.4, "total = 100×30% + 80×30% + 80×40% = 30 + 24 + 32 = 86 分", sz=15, color=WHITE, font="Consolas")
txt(s, 1.1, 6.55, 10, 0.3, "💡 此模式通常用於 ETF 或新上市股票", sz=12, color=DIM)

# ========== Slide 8: 技術面總覽 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "Step 2C  技術面評分 (technical_score)", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)
txt(s, 0.8, 1.15, 11, 0.4, "tech_score = MA×20% + KD×15% + MACD×20% + RSI×15% + 布林×15% + 量能×15%", sz=15, color=WHITE, font="Consolas")
txt(s, 0.8, 1.5, 11, 0.3, "需要至少 20 天日K資料（完整計算需 120 天）", sz=13, color=DIM)

# MA
card(s, 0.8, 2.0, 3.8, 2.3, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 1.0, 2.05, 3.4, 0.35, "📈 均線排列 (20%)", sz=16, color=GREEN, bold=True)
txt(s, 1.0, 2.4, 3.4, 0.25, "比較：MA5>MA10>MA20>MA60>MA120", sz=10, color=WHITE, font="Consolas")
txt(s, 1.0, 2.65, 3.4, 0.25, "依資料量自動調整比較對數", sz=10, color=DIM)
txt(s, 1.0, 2.9, 3.4, 0.3, "score = 符合對數/總對數 × 100", sz=11, color=GOLD, bold=True, font="Consolas")
txt(s, 1.0, 3.25, 3.4, 0.7, "例：4 對中 3 對符合\nscore = 3÷4 × 100 = 75", sz=11, color=DIM)

# KD
card(s, 4.9, 2.0, 3.8, 2.3, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 5.1, 2.05, 3.4, 0.35, "📉 KD 指標 (15%)", sz=16, color=BLUE, bold=True)
txt(s, 5.1, 2.4, 3.4, 0.2, "RSV = (收盤-9日低)÷(9日高-9日低)×100", sz=9, color=WHITE, font="Consolas")
txt(s, 5.1, 2.6, 3.4, 0.2, "K = 2/3×前日K + 1/3×RSV", sz=9, color=WHITE, font="Consolas")
txt(s, 5.1, 2.8, 3.4, 0.2, "D = 2/3×前日D + 1/3×K", sz=9, color=WHITE, font="Consolas")

score_table(s, 5.1, 3.05, 3.4, [
    ("低檔黃金交叉 K<30", "100 分", GREEN),
    ("黃金交叉 K>30", "80 分", BLUE),
    ("低檔整理 K,D在20-40", "70 分", TEAL),
    ("正常區 20-80", "50 分", GOLD),
    ("超買區 >80", "30 分", RED),
])

# MACD
card(s, 9, 2.0, 3.8, 2.3, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 9.2, 2.05, 3.4, 0.35, "📊 MACD (20%)", sz=16, color=PURPLE, bold=True)
txt(s, 9.2, 2.4, 3.4, 0.2, "DIF = EMA(12) - EMA(26)", sz=9, color=WHITE, font="Consolas")
txt(s, 9.2, 2.6, 3.4, 0.2, "Signal = EMA(DIF, 9)", sz=9, color=WHITE, font="Consolas")
txt(s, 9.2, 2.8, 3.4, 0.2, "Histogram = DIF - Signal", sz=9, color=WHITE, font="Consolas")

score_table(s, 9.2, 3.05, 3.4, [
    ("DIF 上穿 Signal", "100 分", GREEN),
    ("柱狀圖 > 0", "70 分", BLUE),
    ("柱狀圖 < 0", "30 分", RED),
])
txt(s, 9.2, 4.05, 3.4, 0.2, "💡 資料<35天→中性50分", sz=10, color=DIM)

# RSI
card(s, 0.8, 4.6, 3.8, 2.5, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 1.0, 4.65, 3.4, 0.35, "⚡ RSI 強弱指標 (15%)", sz=16, color=ORANGE, bold=True)
txt(s, 1.0, 5.0, 3.4, 0.2, "RSI(14) = 100 - 100/(1+RS)", sz=10, color=WHITE, font="Consolas")
txt(s, 1.0, 5.2, 3.4, 0.2, "RS = 14日漲幅均值÷14日跌幅均值", sz=9, color=WHITE, font="Consolas")

score_table(s, 1.0, 5.5, 3.4, [
    ("50~70 健康多頭", "100 分", GREEN),
    ("40~50 溫和", "80 分", BLUE),
    ("30~40 偏弱", "60 分", TEAL),
    ("> 80 超買", "20 分", RED),
    ("< 30 超賣", "40 分", ORANGE),
])

# Bollinger Bands
card(s, 4.9, 4.6, 3.8, 2.5, fill=RGBColor(0x2A, 0x3A, 0x15))
txt(s, 5.1, 4.65, 3.4, 0.35, "🎯 布林通道 (15%)", sz=16, color=TEAL, bold=True)
txt(s, 5.1, 5.0, 3.4, 0.2, "中軌 = SMA(20)", sz=9, color=WHITE, font="Consolas")
txt(s, 5.1, 5.18, 3.4, 0.2, "上軌 = 中軌 + 2×標準差", sz=9, color=WHITE, font="Consolas")
txt(s, 5.1, 5.36, 3.4, 0.2, "%B = (收盤-下軌)÷(上軌-下軌)", sz=9, color=GOLD, font="Consolas")

score_table(s, 5.1, 5.6, 3.4, [
    ("%B >= 1.0 (破上軌)", "70 分", TEAL),
    ("%B >= 0.8 (強勢)", "100 分", GREEN),
    ("%B >= 0.5 (偏多)", "80 分", BLUE),
    ("%B >= 0.2 (偏空)", "40 分", ORANGE),
    ("%B >= 0.0 (弱勢)", "20 分", RED),
    ("%B < 0 (破下軌)", "30 分", ORANGE),
])

# Volume
card(s, 9, 4.6, 3.8, 2.5, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 9.2, 4.65, 3.4, 0.35, "🔊 量能指標 (15%)", sz=16, color=GOLD, bold=True)
txt(s, 9.2, 5.0, 3.4, 0.2, "量比 = 當日量 ÷ MA20量", sz=10, color=WHITE, font="Consolas")

score_table(s, 9.2, 5.3, 3.4, [
    ("量比 >= 1.5 (放量)", "100 分", GREEN),
    ("量比 >= 1.2 (溫和放量)", "80 分", BLUE),
    ("量比 >= 1.0 (正常)", "60 分", TEAL),
    ("量比 < 1.0 (縮量)", "40 分", ORANGE),
])

txt(s, 9.2, 6.6, 3.4, 0.3, "💡 MA20量 = 20日平均成交量", sz=10, color=DIM)

# ========== Slide 9: 綜合評分計算 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "Step 3  綜合分數計算", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)

# Formula card
card(s, 1.5, 1.5, 10.3, 1.5, fill=RGBColor(0x2A, 0x2A, 0x45))
txt(s, 1.8, 1.6, 9.7, 0.5, "計算公式", sz=16, color=DIM, align=PP_ALIGN.CENTER)
txt(s, 1.8, 2.1, 9.7, 0.7, "total = chip × 40% + fundamental × 35% + technical × 25%", sz=24, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font="Consolas")

# Example
card(s, 1, 3.4, 11.3, 3.3, fill=RGBColor(0x20, 0x20, 0x38))
txt(s, 1.3, 3.5, 10.7, 0.5, "📝  範例：台積電 (2330)", sz=22, color=GOLD, bold=True)

for i, (name, score_val, weight, pct, color) in enumerate([
    ("籌碼面", "82", "× 40%", "= 32.8", GOLD),
    ("基本面", "75", "× 35%", "= 26.3", GREEN),
    ("技術面", "68", "× 25%", "= 17.0", BLUE),
]):
    x = 1.5 + i * 3.5
    card(s, x, 4.1, 3, 1.2, fill=RGBColor(0x30, 0x30, 0x50))
    txt(s, x + 0.1, 4.15, 2.8, 0.4, name, sz=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, 4.55, 1.2, 0.5, score_val, sz=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    txt(s, x + 1.3, 4.6, 0.8, 0.4, weight, sz=14, color=DIM, align=PP_ALIGN.CENTER)
    txt(s, x + 2, 4.6, 0.8, 0.4, pct, sz=14, color=color, bold=True, align=PP_ALIGN.CENTER)

card(s, 4, 5.6, 5.3, 0.9, fill=GOLD)
txt(s, 4.3, 5.65, 4.7, 0.8, "綜合分數 = 32.8 + 26.3 + 17.0 = 76.1", sz=22, color=BG, bold=True, align=PP_ALIGN.CENTER, font="Consolas")

# ========== Slide 10: 右側買法信號 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "右側買法：6 個進場信號", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)
txt(s, 0.8, 1.15, 11, 0.35, "score = 觸發信號的權重加總（滿分 100）", sz=15, color=WHITE, font="Consolas")

signals = [
    ("🔊", "量價齊揚", "收盤 > 前日收盤\nAND 量 >= 1.5×MA20量", 25, GOLD),
    ("🚀", "突破20日高", "收盤 > 前20日最高點", 20, GREEN),
    ("📊", "MACD黃金交叉", "DIF 從下方穿越 Signal\n(前日DIF<=Signal 今日>)", 20, BLUE),
    ("📈", "站回MA20", "前日收盤 < MA20\n今日收盤 >= MA20", 15, PURPLE),
    ("⚡", "KD低檔交叉", "K 上穿 D\nAND K < 30", 12, ORANGE),
    ("🎯", "突破布林上軌", "前日收盤 <= 上軌\n今日收盤 > 上軌", 8, TEAL),
]

for i, (icon, name, desc, weight, color) in enumerate(signals):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.7 + row * 2.6
    card(s, x, y, 3.8, 2.2)
    txt(s, x + 0.1, y + 0.1, 0.6, 0.5, icon, sz=24)
    txt(s, x + 0.7, y + 0.1, 2.2, 0.4, name, sz=17, color=color, bold=True)
    txt(s, x + 0.7, y + 0.5, 2.5, 0.9, desc, sz=11, color=GRAY, font="Consolas")
    # Weight circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 3), Inches(y + 0.15), Inches(0.6), Inches(0.6))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = str(weight)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

txt(s, 0.8, 6.7, 11.5, 0.35, "💡 信號觸發判斷使用最近 120 天日K資料  |  權重合計 = 25+20+20+15+12+8 = 100", sz=13, color=DIM, align=PP_ALIGN.CENTER)

# ========== Slide 10.5: 進階條件篩選 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "右側買法：進階條件篩選", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)
txt(s, 0.8, 1.2, 11, 0.35, "在 6 個基本信號之外，系統額外計算 4 個進階條件", sz=16, color=DIM)

# Condition 1: Today Breakout
card(s, 0.8, 1.8, 5.7, 2.5, fill=RGBColor(0x1E, 0x3A, 0x1E))
txt(s, 1.0, 1.9, 5.3, 0.45, "🚀  今日突破 (today_breakout)", sz=18, color=GREEN, bold=True)
txt(s, 1.0, 2.4, 5.3, 0.3, "判斷條件：", sz=13, color=DIM)
txt(s, 1.0, 2.7, 5.3, 0.3, "量價齊揚 = True  AND  突破20日高點 = True", sz=13, color=WHITE, font="Consolas")
txt(s, 1.0, 3.1, 5.3, 0.6, "意義：同時出現量增價揚與突破高點\n代表多方力道強勁的突破訊號", sz=12, color=GRAY)

# Condition 2: Weekly Trend Up
card(s, 7, 1.8, 5.7, 2.5, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 7.2, 1.9, 5.3, 0.45, "📈  週趨勢向上 (weekly_trend_up)", sz=18, color=BLUE, bold=True)
txt(s, 7.2, 2.4, 5.3, 0.3, "判斷條件：", sz=13, color=DIM)
txt(s, 7.2, 2.7, 5.3, 0.3, "MA5 > MA20  AND  MA5(今) > MA5(3日前)", sz=13, color=WHITE, font="Consolas")
txt(s, 7.2, 3.1, 5.3, 0.6, "意義：短期均線在中期均線上方\n且短期趨勢持續上升中", sz=12, color=GRAY)

# Condition 3: Risk Level
card(s, 0.8, 4.6, 5.7, 2.5, fill=RGBColor(0x3A, 0x2E, 0x15))
txt(s, 1.0, 4.7, 5.3, 0.45, "⚠️  風險等級 (risk_level)", sz=18, color=ORANGE, bold=True)
txt(s, 1.0, 5.2, 5.3, 0.3, "計算公式：20日報酬率標準差 × √252", sz=13, color=WHITE, font="Consolas")

score_table(s, 1.0, 5.6, 5.3, [
    ("波動率<25% 且 分數>=60", "low 低風險", GREEN),
    ("波動率<40% 或 分數>=45", "medium 中風險", GOLD),
    ("其他", "high 高風險", RED),
])

# Condition 4: Strong Recommend
card(s, 7, 4.6, 5.7, 2.5, fill=RGBColor(0x2A, 0x1E, 0x3A))
txt(s, 7.2, 4.7, 5.3, 0.45, "🌟  強力推薦 (strong_recommend)", sz=18, color=PURPLE, bold=True)
txt(s, 7.2, 5.2, 5.3, 0.3, "必須同時滿足 4 個條件：", sz=13, color=DIM)
txt(s, 7.2, 5.55, 5.3, 0.25, "① 信號加權分數 >= 60", sz=12, color=WHITE, font="Consolas")
txt(s, 7.2, 5.8, 5.3, 0.25, "② 觸發信號數 >= 3", sz=12, color=WHITE, font="Consolas")
txt(s, 7.2, 6.05, 5.3, 0.25, "③ 週趨勢向上 = True", sz=12, color=WHITE, font="Consolas")
txt(s, 7.2, 6.3, 5.3, 0.25, "④ 風險等級 ≠ high", sz=12, color=WHITE, font="Consolas")

# ========== Slide 11: 買賣點預測 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "買賣點預測公式", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)

# Entry
card(s, 0.8, 1.5, 3.8, 3.0, fill=RGBColor(0x1E, 0x3A, 0x1E))
txt(s, 1, 1.6, 3.4, 0.5, "🟢  進場價", sz=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1, 2.2, 3.4, 0.4, "entry = 最新收盤價", sz=16, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
txt(s, 1, 2.8, 3.4, 0.5, "直接參考當下價位\n信號確認即可進場", sz=13, color=DIM, align=PP_ALIGN.CENTER)

# Stop loss
card(s, 4.9, 1.5, 3.8, 3.0, fill=RGBColor(0x3A, 0x1E, 0x1E))
txt(s, 5.1, 1.6, 3.4, 0.5, "🔴  停損價", sz=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
txt(s, 5.1, 2.15, 3.4, 0.3, "stop = max(MA20, 20日最低)", sz=13, color=WHITE, font="Consolas")
txt(s, 5.1, 2.45, 3.4, 0.3, "若 stop >= entry:", sz=12, color=DIM, font="Consolas")
txt(s, 5.1, 2.7, 3.4, 0.3, "  stop = entry × 0.95", sz=13, color=ORANGE, font="Consolas")
txt(s, 5.1, 3.1, 3.4, 0.5, "取最近支撐位\n確保停損低於進場價", sz=13, color=DIM, align=PP_ALIGN.CENTER)

# Target
card(s, 9, 1.5, 3.8, 3.0, fill=RGBColor(0x1E, 0x2E, 0x4A))
txt(s, 9.2, 1.6, 3.4, 0.5, "🎯  目標價", sz=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 9.2, 2.15, 3.4, 0.3, "risk = entry - stop_loss", sz=13, color=WHITE, font="Consolas")
txt(s, 9.2, 2.45, 3.4, 0.3, "target = entry + M × risk", sz=13, color=GOLD, bold=True, font="Consolas")
txt(s, 9.2, 2.75, 3.4, 0.3, "M: ≥60→2.0  ≥35→1.5  else→1.0", sz=11, color=ORANGE, font="Consolas")
txt(s, 9.2, 3.05, 3.4, 0.3, "若 risk<=0: entry × 1.05", sz=12, color=DIM, font="Consolas")

# Action suggestions
txt(s, 0.8, 4.8, 11, 0.5, "動作建議（依信號分數判定）", sz=22, color=GOLD, bold=True)

card(s, 0.8, 5.4, 3.8, 1.5, fill=RGBColor(0x1E, 0x3A, 0x1E))
txt(s, 1, 5.5, 3.4, 0.45, "✅  BUY 建議買入", sz=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1, 5.95, 3.4, 0.4, "score >= 60", sz=16, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

card(s, 4.9, 5.4, 3.8, 1.5, fill=RGBColor(0x3A, 0x3A, 0x1E))
txt(s, 5.1, 5.5, 3.4, 0.45, "⏸️  HOLD 觀望等待", sz=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 5.1, 5.95, 3.4, 0.4, "35 <= score < 60", sz=16, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

card(s, 9, 5.4, 3.8, 1.5, fill=RGBColor(0x3A, 0x1E, 0x1E))
txt(s, 9.2, 5.5, 3.4, 0.45, "❌  AVOID 暫不建議", sz=20, color=RED, bold=True, align=PP_ALIGN.CENTER)
txt(s, 9.2, 5.95, 3.4, 0.4, "score < 35", sz=16, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

# ========== Slide 12: 每日自動流程 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 0.8, 0.4, 11, 0.7, "每日自動更新流程", sz=34, color=GOLD, bold=True)
bar(s, 0.8, 1.0, 3)

times = [
    ("Step 1", "資料收集", "全市場收盤價\n法人、融資融券\n營收、財報", BLUE),
    ("Step 2", "候選股篩選", "量能異常偵測\n+ Top 500 保底\n→ ~500 檔候選", GREEN),
    ("Step 3", "三因子評分", "籌碼+基本面\n+技術面\n加權排名", GOLD),
    ("完成", "結果上線", "Dashboard\n顯示最新\n排名結果", TEAL),
]

for i, (time, title, desc, color) in enumerate(times):
    x = 0.8 + i * 3.1
    card(s, x + 0.3, 1.5, 1.8, 0.5, fill=color)
    txt(s, x + 0.3, 1.5, 1.8, 0.5, time, sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    if i < 3:
        bar(s, x + 2.1, 1.7, 1.5, DIM)

    card(s, x + 0.05, 2.3, 2.6, 2.5)
    txt(s, x + 0.15, 2.4, 2.4, 0.4, title, sz=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + 0.15, 2.8, 2.4, 1.2, desc, sz=13, color=GRAY, align=PP_ALIGN.CENTER)

card(s, 1, 5.5, 11.3, 1.2, fill=RGBColor(0x2A, 0x2A, 0x45))
txt(s, 1.3, 5.6, 10.7, 0.5, "💡  全自動運作，無需手動操作", sz=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.3, 6.1, 10.7, 0.4, "系統自動偵測交易日，非交易日自動略過  |  排程時間可在設定頁面調整  |  AI 報告為手動生成", sz=14, color=DIM, align=PP_ALIGN.CENTER)

# ========== Slide 13: 結尾 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(s)
txt(s, 1, 2, 11, 1, "讓數據說話", sz=48, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
bar(s, 4.5, 3.2, 4)
txt(s, 1, 3.5, 11, 0.8, "台股智慧選股系統 v2.2", sz=24, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 1, 4.8, 11, 0.5, "三因子量化評分  ×  AI 智慧分析  ×  動能信號偵測", sz=18, color=DIM, align=PP_ALIGN.CENTER)

# Save
out = "/Users/weihung/Desktop/project/stock-system/docs/台股智慧選股系統-指標公式與評分邏輯.pptx"
prs.save(out)
print(f"Generated: {out} ({len(prs.slides)} slides)")
